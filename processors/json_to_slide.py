import json
import logging
import re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# === Маппинг названий фигур из парсера в python-pptx ===
# Парсер пишет "ROUNDED_RECTANGLE (5)", нам нужен MSO_SHAPE.ROUNDED_RECTANGLE
SHAPE_TYPE_MAP: dict[str, MSO_SHAPE] = {
    "ROUNDED_RECTANGLE": MSO_SHAPE.ROUNDED_RECTANGLE,
    "RECTANGLE": MSO_SHAPE.RECTANGLE,
    "DIAMOND": MSO_SHAPE.DIAMOND,
    "OVAL": MSO_SHAPE.OVAL,
    "ROUND_2_DIAG_RECTANGLE": MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
    "ROUND_1_RECTANGLE": MSO_SHAPE.ROUND_1_RECTANGLE,
    "ROUND_2_SAME_RECTANGLE": MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
    "ISOSCELES_TRIANGLE": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "RIGHT_TRIANGLE": MSO_SHAPE.RIGHT_TRIANGLE,
    "PARALLELOGRAM": MSO_SHAPE.PARALLELOGRAM,
    "TRAPEZOID": MSO_SHAPE.TRAPEZOID,
    "HEXAGON": MSO_SHAPE.HEXAGON,
    "OCTAGON": MSO_SHAPE.OCTAGON,
    "CROSS": MSO_SHAPE.CROSS,
    "STAR_5_POINT": MSO_SHAPE.STAR_5_POINT,
    "RIGHT_ARROW": MSO_SHAPE.RIGHT_ARROW,
    "LEFT_ARROW": MSO_SHAPE.LEFT_ARROW,
    "CHEVRON": MSO_SHAPE.CHEVRON,
    "PENTAGON": MSO_SHAPE.PENTAGON,
    "HEART": MSO_SHAPE.HEART,
    "CLOUD": MSO_SHAPE.CLOUD,
    "LIGHTNING_BOLT": MSO_SHAPE.LIGHTNING_BOLT,
}

ALIGNMENT_MAP: dict[str, PP_ALIGN] = {
    "LEFT": PP_ALIGN.LEFT,
    "CENTER": PP_ALIGN.CENTER,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
}

ANCHOR_MAP: dict[str, str] = {
    "top": "t",
    "middle": "ctr",
    "bottom": "b",
}

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Преобразует '#RRGGBB' в RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def parse_enum_name(raw: str) -> str:
    """
    Извлекает имя из строки вида 'ROUNDED_RECTANGLE (5)'.
    Возвращает 'ROUNDED_RECTANGLE'.
    """
    if not raw:
        return ""
    match = re.match(r"^([A-Z_0-9]+)", raw)
    return match.group(1) if match else raw.upper().replace(" ", "_")


def resolve_color(
    color_source: str | None,
    color_hex: str | None,
    theme_colors: dict[str, str],
    theme_color_name: str | None = None,
) -> str | None:
    """
    Определяет итоговый hex цвет.
    Приоритет: прямой hex → theme_color → inherited (None).
    """
    if color_hex:
        return color_hex
    if theme_color_name and theme_color_name in theme_colors:
        return theme_colors[theme_color_name]
    return None


def apply_solid_fill(shape, fill_data: dict, theme_colors: dict) -> None:
    """Применяет сплошную заливку к фигуре."""
    hex_color = fill_data.get("rgb")
    if not hex_color and fill_data.get("theme_color"):
        hex_color = theme_colors.get(fill_data["theme_color"])
    if hex_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(hex_color)


def apply_gradient_fill(shape, fill_data: dict) -> None:
    """
    Применяет градиентную заливку через XML.
    python-pptx не умеет gradient — делаем руками.
    """
    sp_pr = shape._element.find(f".//{{{A_NS}}}spPr")
    if sp_pr is None:
        return

    # Удаляем существующие заливки
    for old in sp_pr.findall(f"{{{A_NS}}}solidFill") + sp_pr.findall(f"{{{A_NS}}}gradFill") + sp_pr.findall(f"{{{A_NS}}}noFill"):
        sp_pr.remove(old)

    grad_fill = etree.SubElement(sp_pr, f"{{{A_NS}}}gradFill")
    gs_lst = etree.SubElement(grad_fill, f"{{{A_NS}}}gsLst")

    for stop in fill_data.get("gradient_stops", []):
        pos_val = str(int(stop["position"] * 100000))
        gs = etree.SubElement(gs_lst, f"{{{A_NS}}}gs")
        gs.set("pos", pos_val)
        srgb = etree.SubElement(gs, f"{{{A_NS}}}srgbClr")
        srgb.set("val", stop["color"].lstrip("#"))

    lin = etree.SubElement(grad_fill, f"{{{A_NS}}}lin")
    angle_val = str(int(fill_data.get("gradient_angle", 0) * 60000))
    lin.set("ang", angle_val)
    lin.set("scaled", "0")


def apply_no_line(shape) -> None:
    """Убирает обводку у фигуры."""
    shape.line.fill.background()


def apply_line(shape, line_data: dict) -> None:
    """Применяет настройки линии (обводки) к фигуре."""
    if line_data.get("border") == "none":
        apply_no_line(shape)
        return
    if line_data.get("color"):
        shape.line.color.rgb = hex_to_rgb(line_data["color"])
    if line_data.get("width_pt"):
        shape.line.width = Pt(line_data["width_pt"])


def apply_shadow(shape, shadow_data: dict) -> None:
    """Применяет тень к фигуре через XML."""
    sp_pr = shape._element.find(f".//{{{A_NS}}}spPr")
    if sp_pr is None:
        return

    effect_lst = sp_pr.find(f"{{{A_NS}}}effectLst")
    if effect_lst is None:
        effect_lst = etree.SubElement(sp_pr, f"{{{A_NS}}}effectLst")

    outer_shdw = etree.SubElement(effect_lst, f"{{{A_NS}}}outerShdw")
    outer_shdw.set("blurRad", str(shadow_data.get("blur_radius_emu", 76200)))
    outer_shdw.set("dist", str(shadow_data.get("distance_emu", 38100)))
    outer_shdw.set("dir", str(shadow_data.get("direction_emu", 5400000)))
    outer_shdw.set("algn", "tl")
    outer_shdw.set("rotWithShape", "0")

    color = shadow_data.get("color", "#000000").lstrip("#")
    alpha = shadow_data.get("alpha", 15000)

    srgb_clr = etree.SubElement(outer_shdw, f"{{{A_NS}}}srgbClr")
    srgb_clr.set("val", color)
    alpha_el = etree.SubElement(srgb_clr, f"{{{A_NS}}}alpha")
    alpha_el.set("val", str(alpha))


def apply_adjustments(shape, adjustments: list[float]) -> None:
    """Устанавливает adjustment values (закругления и т.д.)."""
    for i, val in enumerate(adjustments):
        try:
            shape.adjustments[i] = val
        except IndexError:
            logger.warning(f"Adjustment index {i} out of range for shape")


def apply_text_frame(shape, tf_data: dict, theme_colors: dict) -> None:
    """Заполняет текстовый фрейм фигуры из parsed JSON."""
    tf = shape.text_frame

    if tf_data.get("word_wrap") is not None:
        tf.word_wrap = tf_data["word_wrap"]

    auto_size_name = parse_enum_name(str(tf_data.get("auto_size", "NONE")))
    auto_size_map = {
        "NONE": MSO_AUTO_SIZE.NONE,
        "SHAPE_TO_FIT_TEXT": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        "TEXT_TO_FIT_SHAPE": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    }
    if auto_size_name in auto_size_map:
        tf.auto_size = auto_size_map[auto_size_name]

    vert = tf_data.get("vertical_alignment")
    if vert and vert in ANCHOR_MAP:
        tf._txBody.bodyPr.set("anchor", ANCHOR_MAP[vert])

    paragraphs = tf_data.get("paragraphs", [])
    for p_idx, p_data in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        align_raw = parse_enum_name(p_data.get("alignment", ""))
        if align_raw in ALIGNMENT_MAP:
            p.alignment = ALIGNMENT_MAP[align_raw]

        for run_data in p_data.get("runs", []):
            run = p.add_run()
            run.text = run_data.get("text", "")

            font = run.font
            if run_data.get("font_name"):
                font.name = run_data["font_name"]
            if run_data.get("font_size_pt"):
                font.size = Pt(run_data["font_size_pt"])
            if run_data.get("bold"):
                font.bold = True
            if run_data.get("italic"):
                font.italic = True

            run_color = resolve_color(
                run_data.get("color_source"),
                run_data.get("color"),
                theme_colors,
                run_data.get("theme_color"),
            )
            if run_color:
                font.color.rgb = hex_to_rgb(run_color)


def add_auto_shape(slide, shape_data: dict, theme_colors: dict) -> None:
    """Добавляет AutoShape (прямоугольник, ромб и т.д.) на слайд."""
    raw_type = parse_enum_name(shape_data.get("auto_shape_type", "RECTANGLE"))
    mso_type = SHAPE_TYPE_MAP.get(raw_type)

    if not mso_type:
        logger.warning(f"Unknown auto_shape_type: {raw_type}, skipping shape {shape_data.get('name')}")
        return

    shape = slide.shapes.add_shape(
        mso_type,
        Inches(shape_data["left_inches"]),
        Inches(shape_data["top_inches"]),
        Inches(shape_data["width_inches"]),
        Inches(shape_data["height_inches"]),
    )

    if shape_data.get("rotation"):
        shape.rotation = shape_data["rotation"]

    fill = shape_data.get("fill", {})
    fill_type = fill.get("type", "")

    if fill_type == "solid":
        apply_solid_fill(shape, fill, theme_colors)
    elif fill_type == "gradient":
        apply_gradient_fill(shape, fill)
    elif "BACKGROUND" in fill_type:
        shape.fill.background()

    if shape_data.get("line"):
        apply_line(shape, shape_data["line"])

    if shape_data.get("shadow"):
        apply_shadow(shape, shape_data["shadow"])

    if shape_data.get("adjustments"):
        apply_adjustments(shape, shape_data["adjustments"])

    if shape_data.get("text_frame"):
        apply_text_frame(shape, shape_data["text_frame"], theme_colors)

    logger.debug(f"Added auto_shape: {shape_data.get('name')}")


def add_text_box(slide, shape_data: dict, theme_colors: dict) -> None:
    """Добавляет TextBox на слайд."""
    txBox = slide.shapes.add_textbox(
        Inches(shape_data["left_inches"]),
        Inches(shape_data["top_inches"]),
        Inches(shape_data["width_inches"]),
        Inches(shape_data["height_inches"]),
    )

    if shape_data.get("rotation"):
        txBox.rotation = shape_data["rotation"]

    fill = shape_data.get("fill", {})
    fill_type = fill.get("type", "")

    if fill_type == "solid":
        apply_solid_fill(txBox, fill, theme_colors)
    elif "BACKGROUND" in fill_type:
        txBox.fill.background()

    if shape_data.get("text_frame"):
        apply_text_frame(txBox, shape_data["text_frame"], theme_colors)

    logger.debug(f"Added text_box: {shape_data.get('name')}")


def add_freeform_placeholder(slide, shape_data: dict, theme_colors: dict) -> None:
    """
    Freeform (иконки) — у нас нет path data, поэтому ставим
    прямоугольник-заглушку того же размера. Лучше чем ничего.
    """
    logger.debug(f"Skipped freeform: {shape_data.get('name')}")
    return


def process_shape(slide, shape_data: dict, theme_colors: dict) -> None:
    """
    Роутер: смотрит на shape_type и вызывает нужную функцию.
    """
    raw_shape_type = parse_enum_name(shape_data.get("shape_type", ""))

    if raw_shape_type == "GROUP":
        for child in shape_data.get("children", []):
            process_shape(slide, child, theme_colors)
        return

    if raw_shape_type == "AUTO_SHAPE":
        add_auto_shape(slide, shape_data, theme_colors)
    elif raw_shape_type == "TEXT_BOX":
        add_text_box(slide, shape_data, theme_colors)
    elif raw_shape_type == "FREEFORM":
        add_freeform_placeholder(slide, shape_data, theme_colors)
    elif raw_shape_type == "PICTURE":
        logger.info(f"Skipping picture: {shape_data.get('name')} (image rebuild not implemented)")
    elif raw_shape_type == "TABLE":
        logger.info(f"Skipping table: {shape_data.get('name')} (table rebuild not implemented)")
    elif raw_shape_type == "CHART":
        logger.info(f"Skipping chart: {shape_data.get('name')} (chart rebuild not implemented)")
    else:
        logger.warning(f"Unknown shape_type: {raw_shape_type}, skipping: {shape_data.get('name')}")


def apply_slide_background(slide, bg_data: dict | None, theme_colors: dict) -> None:
    """Устанавливает фон слайда."""
    if not bg_data:
        return

    bg_type = bg_data.get("type", "")

    if bg_type == "solid":
        hex_color = bg_data.get("rgb")
        if not hex_color and bg_data.get("theme_color"):
            hex_color = theme_colors.get(bg_data["theme_color"])
        if hex_color:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(hex_color)


def build_slide_from_json(json_path: Path, output_path: Path) -> None:
    """
    Собирает презентацию из parsed JSON файла.

    :param json_path: Путь к parsed.json.
    :param output_path: Путь для сохранения .pptx.
    """
    data = json.loads(json_path.read_text(encoding="utf-8"))
    prs = Presentation()

    prs.slide_width = Inches(data.get("width_inches", 13.3333))
    prs.slide_height = Inches(data.get("height_inches", 7.5))

    theme_colors = data.get("theme_colors", {})
    blank_layout = prs.slide_layouts[6]

    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        apply_slide_background(slide, slide_data.get("background"), theme_colors)

        for shape_data in slide_data.get("shapes", []):
            try:
                process_shape(slide, shape_data, theme_colors)
            except Exception as e:
                logger.error(f"Error processing shape '{shape_data.get('name')}': {e}")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Presentation saved to {output_path}")