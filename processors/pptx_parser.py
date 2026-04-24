import json
import logging
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

logger = logging.getLogger(__name__)

def extract_theme_colors(pptx_path: Path) -> Dict[str, str]:
    """Extracts theme colors from pptx file."""
    try:
        with zipfile.ZipFile(str(pptx_path)) as z:
            theme_xml = z.read("ppt/theme/theme1.xml")
        tree = etree.fromstring(theme_xml)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = tree.find(".//a:clrScheme", ns)
        if clr_scheme is None:
            return {}
        
        theme_map = {}
        for child in clr_scheme:
            color_name = child.tag.split("}")[-1]
            srgb = child.find("a:srgbClr", ns)
            sys_clr = child.find("a:sysClr", ns)
            if srgb is not None:
                theme_map[color_name] = f"#{srgb.get('val')}"
            elif sys_clr is not None:
                theme_map[color_name] = f"#{sys_clr.get('lastClr')}"
        
        alias_map = {
            "bg1": "lt1", "bg2": "lt2",
            "tx1": "dk1", "tx2": "dk2"
        }
        for alias, real in alias_map.items():
            if real in theme_map:
                theme_map[alias] = theme_map[real]
        return theme_map
    except Exception as e:
        logger.error(f"Error extracting theme colors: {e}")
        return {}

def emu_to_inches(emu: Any) -> float:
    try:
        return round(int(emu) / 914400, 4)
    except:
        return 0.0

def emu_to_pt(emu: Any) -> float:
    try:
        return round(int(emu) / 12700, 2)
    except:
        return 0.0

def color_to_hex(color: Any, shape: Any = None) -> Optional[Dict[str, Any]]:
    try:
        if color and hasattr(color, 'rgb') and color.rgb:
            return {"rgb": f"#{str(color.rgb)}"}
    except:
        pass
    
    if shape is not None:
        try:
            sp = shape._element
            solid_fill = sp.find(".//" + qn("a:solidFill"))
            if solid_fill is not None:
                scheme_clr = solid_fill.find(qn("a:schemeClr"))
                if scheme_clr is not None:
                    return {"theme_color": scheme_clr.get("val")}
        except:
            pass
    return None

def get_fill_data(fill: Any, shape: Any = None) -> Optional[Dict[str, Any]]:
    try:
        if fill.type is None:
            return None
        if "SOLID" in str(fill.type).upper(): # solid
            fill_data = {"type": "solid"}
            if color := color_to_hex(fill.fore_color, shape):
                fill_data.update(color)
            
            # Theme colors fallback
            try:
                sp = shape._element
                solid_fill = sp.find(".//" + qn("a:solidFill"))
                if solid_fill is not None:
                    scheme_clr = solid_fill.find(qn("a:schemeClr"))
                    if scheme_clr is not None:
                        fill_data["theme_color"] = scheme_clr.get("val")
            except:
                pass
            return fill_data
        elif "GRADIENT" in str(fill.type).upper(): # gradient
            fill_data = {"type": "gradient"}
            try:
                spPr = shape._element.find(qn("p:spPr"))
                if spPr is None:
                    spPr = shape._element.find(qn("a:spPr"))
                grad_fill = None
                if spPr is not None:
                    grad_fill = spPr.find(qn("a:gradFill"))
                if grad_fill is None:
                    grad_fill = shape._element.find(".//" + qn("a:gradFill"))
                
                if grad_fill is not None:
                    lin = grad_fill.find(qn("a:lin"))
                    if lin is not None:
                        ang = lin.get("ang")
                        if ang:
                            fill_data["gradient_angle"] = int(ang) / 60000
                    
                    gs_lst = grad_fill.find(qn("a:gsLst"))
                    if gs_lst is not None:
                        stops = []
                        for gs in gs_lst.findall(qn("a:gs")):
                            stop = {"position": int(gs.get("pos", "0")) / 100000}
                            srgb = gs.find(qn("a:srgbClr"))
                            if srgb is not None:
                                stop["color"] = f"#{srgb.get('val')}"
                            else:
                                scheme = gs.find(qn("a:schemeClr"))
                                if scheme is not None:
                                    stop["theme_color"] = scheme.get("val")
                                    lum_mod = scheme.find(qn("a:lumMod"))
                                    if lum_mod is not None:
                                        stop["lum_mod"] = int(lum_mod.get("val")) / 1000
                                    lum_off = scheme.find(qn("a:lumOff"))
                                    if lum_off is not None:
                                        stop["lum_off"] = int(lum_off.get("val")) / 1000
                            stops.append(stop)
                        fill_data["gradient_stops"] = stops
            except Exception:
                pass
            
            # Theme colors logic
            try:
                sp = shape._element
                solid_fill = sp.find(".//" + qn("a:solidFill"))
                if solid_fill is not None:
                    scheme_clr = solid_fill.find(qn("a:schemeClr"))
                    if scheme_clr is not None:
                        fill_data["theme_color"] = scheme_clr.get("val")
            except Exception:
                pass
                
            return fill_data
        return {"type": str(fill.type)}
    except:
        return None

def get_line_data(line: Any, shape: Any = None) -> Optional[Dict[str, Any]]:
    try:
        if not line:
            return None
        if hasattr(line, 'fill') and line.fill and hasattr(line.fill, 'type') and "BACKGROUND" in str(line.fill.type):
            return {"border": "none"}
        data = {}
        if color := color_to_hex(line.color, shape):
            data["color"] = color
        if line.width:
            data["width_pt"] = emu_to_pt(line.width)
        if line.dash_style:
            data["dash_style"] = str(line.dash_style)
        return data if data else None
    except:
        return None

def get_shadow_data(shape: Any) -> Optional[Dict[str, Any]]:
    try:
        sp = shape._element
        effect_lst = sp.find(".//" + qn("a:effectLst"))
        if effect_lst is not None:
            outer_shdw = effect_lst.find(qn("a:outerShdw"))
            if outer_shdw is not None:
                data = {
                    "blur_radius_emu": int(outer_shdw.get("blurRad", 0)),
                    "distance_emu": int(outer_shdw.get("dist", 0)),
                    "direction": int(outer_shdw.get("dir", 0)),
                }
                color = outer_shdw.find(qn("a:srgbClr"))
                if color is not None:
                    data["color"] = f"#{color.get('val')}"
                    alpha = color.find(qn("a:alpha"))
                    if alpha is not None:
                        data["alpha"] = int(alpha.get("val"))
                return data
    except:
        pass
    return None

def parse_shape(shape: Any, output_dir: Path, group_offset: Optional[Dict[str, float]] = None) -> Dict[str, Any]:
    shape_data: Dict[str, Any] = {
        "shape_id": getattr(shape, "shape_id", None),
        "name": getattr(shape, "name", None),
        "shape_type": str(shape.shape_type),
        "left_inches": emu_to_inches(shape.left),
        "top_inches": emu_to_inches(shape.top),
        "width_inches": emu_to_inches(shape.width),
        "height_inches": emu_to_inches(shape.height),
        "rotation": getattr(shape, "rotation", 0),
    }

    if group_offset:
        chOff = group_offset["chOff"]
        chExt = group_offset["chExt"]
        
        # Applying formula
        abs_left_emu = group_offset["off_x"] + (shape.left - chOff["x"]) * (group_offset["ext"]["cx"] / chExt["cx"])
        abs_top_emu = group_offset["off_y"] + (shape.top - chOff["y"]) * (group_offset["ext"]["cy"] / chExt["cy"])
        
        shape_data["left_inches"] = emu_to_inches(abs_left_emu)
        shape_data["top_inches"] = emu_to_inches(abs_top_emu)
        shape_data["width_inches"] = emu_to_inches(shape.width * (group_offset["ext"]["cx"] / chExt["cx"]))
        shape_data["height_inches"] = emu_to_inches(shape.height * (group_offset["ext"]["cy"] / chExt["cy"]))

    try:
        if hasattr(shape, "auto_shape_type"):
            shape_data["auto_shape_type"] = str(shape.auto_shape_type)
    except:
        pass

    try:
        if hasattr(shape, "fill"):
            if fill_data := get_fill_data(shape.fill, shape):
                shape_data["fill"] = fill_data
    except:
        pass

    try:
        if hasattr(shape, "line"):
            if line_data := get_line_data(shape.line):
                shape_data["line"] = line_data
    except:
        pass

    if shadow_data := get_shadow_data(shape):
        shape_data["shadow"] = shadow_data

    try:
        if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
            shape_data["adjustments"] = [v for v in shape.adjustments]
    except:
        pass

    if shape.has_text_frame:
        tf = shape.text_frame
        tf_data = {"word_wrap": tf.word_wrap, "auto_size": str(tf.auto_size), "paragraphs": []}
        try:
            anchor = tf._txBody.bodyPr.get("anchor")
            mapping = {"t": "top", "ctr": "middle", "b": "bottom"}
            tf_data["vertical_alignment"] = mapping.get(anchor, anchor)
        except:
            pass
        for p in tf.paragraphs:
            p_data = {}
            if p.alignment is not None: p_data["alignment"] = str(p.alignment)
            if p.level is not None: p_data["level"] = p.level
            if p.space_before is not None: p_data["space_before_pt"] = emu_to_pt(p.space_before)
            if p.space_after is not None: p_data["space_after_pt"] = emu_to_pt(p.space_after)
            p_data["runs"] = []
            for r in p.runs:
                r_data = {"text": r.text}
                if r.font.name: r_data["font_name"] = r.font.name
                if r.font.size: r_data["font_size_pt"] = emu_to_pt(r.font.size)
                if r.font.bold is not None: r_data["bold"] = r.font.bold
                if r.font.italic is not None: r_data["italic"] = r.font.italic
                if r.font.underline is not None: r_data["underline"] = r.font.underline
                try:
                    run_element = r._r
                    solid_fill = run_element.find(".//" + qn("a:solidFill"))
                    if solid_fill is not None:
                        scheme_clr = solid_fill.find(qn("a:schemeClr"))
                        if scheme_clr is not None:
                            r_data["theme_color"] = scheme_clr.get("val")
                    if "color" not in r_data and "theme_color" not in r_data:
                        r_data["color_source"] = "inherited"
                except:
                    pass
                try: r_data["color"] = color_to_hex(r.font.color)
                except: pass
                if r.hyperlink.address: r_data["hyperlink"] = r.hyperlink.address
                p_data["runs"].append(r_data)
            tf_data["paragraphs"].append(p_data)
        shape_data["text_frame"] = tf_data

    if hasattr(shape, "shapes"):
        # Log to validate group coordinate assumptions
        logger.debug(f"Parsing Group: {shape.name}")
        
        # Get xfrm element for group
        grpSpPr = shape._element.find(qn("p:grpSpPr"))
        if grpSpPr is None:
            grpSpPr = shape._element.find(qn("a:grpSpPr"))
        xfrm = grpSpPr.find(qn("a:xfrm")) if grpSpPr is not None else None
        
        if xfrm is not None:
            off = xfrm.find(qn("a:off"))
            ext = xfrm.find(qn("a:ext"))
            chOff = xfrm.find(qn("a:chOff"))
            chExt = xfrm.find(qn("a:chExt"))
            
            if off is not None and ext is not None and chOff is not None and chExt is not None:
                offset = {
                    "off_x": int(off.get("x")),
                    "off_y": int(off.get("y")),
                    "ext": {"cx": int(ext.get("cx")), "cy": int(ext.get("cy"))},
                    "chOff": {"x": int(chOff.get("x")), "y": int(chOff.get("y"))},
                    "chExt": {"cx": int(chExt.get("cx")), "cy": int(chExt.get("cy"))}
                }
            else:
                offset = None
        else:
            offset = None
            
        shape_data["children"] = [parse_shape(child, output_dir, offset) for child in shape.shapes]
        for child in shape.shapes:
            # Get child xfrm element
            c_xfrm = child._element.find(f".//{{{qn('a:xfrm')}}}xfrm")
            if c_xfrm is not None:
                c_off = c_xfrm.find(qn("a:off"))
                c_ext = c_xfrm.find(qn("a:ext"))
                logger.debug(f"Child {child.name}: off=({c_off.get('x')}, {c_off.get('y')}), ext=({c_ext.get('cx')}, {c_ext.get('cy')})")

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or "Picture" in shape.name:
        image_dir = output_dir / "images"
        image_dir.mkdir(parents=True, exist_ok=True)
        filename = f"image_{shape.shape_id}.png"
        with open(image_dir / filename, "wb") as f:
            f.write(shape.image.blob)
        shape_data["image"] = {"filename": filename, "content_type": shape.image.content_type}

    if shape.has_table:
        data = []
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                cell_data = {"text": cell.text_frame.text}
                try: cell_data["fill_color"] = color_to_hex(cell.fill.fore_color)
                except: pass
                row_data.append(cell_data)
            data.append(row_data)
        shape_data["table"] = {"rows": len(shape.table.rows), "cols": len(shape.table.columns), "data": data}

    if shape.has_chart:
        shape_data["chart"] = {"chart_type": str(shape.chart.chart_type), "chart_data": "not_extracted"}

    return shape_data

def parse_presentation(pptx_path: Path, output_dir: Path) -> Dict[str, Any]:
    prs = Presentation(str(pptx_path))
    presentation_data = {
        "width_inches": emu_to_inches(prs.slide_width),
        "height_inches": emu_to_inches(prs.slide_height),
        "theme_colors": extract_theme_colors(pptx_path),
        "slides": []
    }
    for slide in prs.slides:
        slide_data = {"shapes": []}
        try:
            # Check slide background, then layout, then master
            fill = None
            if slide.background.fill.type is not None:
                fill = slide.background.fill
            elif slide.slide_layout.background.fill.type is not None:
                fill = slide.slide_layout.background.fill
            elif slide.slide_master.background.fill.type is not None:
                fill = slide.slide_master.background.fill
            
            if fill and (bg := get_fill_data(fill)):
                slide_data["background"] = bg
        except: pass
        for shape in slide.shapes:
            if shape.shape_type == 13: # CONNECTOR
                c_data = parse_shape(shape, output_dir)
                c_data.update({
                    "begin_x_inches": emu_to_inches(shape.begin_x),
                    "begin_y_inches": emu_to_inches(shape.begin_y),
                    "end_x_inches": emu_to_inches(shape.end_x),
                    "end_y_inches": emu_to_inches(shape.end_y),
                    "line_color": color_to_hex(shape.line.color),
                    "line_width_pt": emu_to_pt(shape.line.width)
                })
                slide_data["shapes"].append(c_data)
            else:
                slide_data["shapes"].append(parse_shape(shape, output_dir))
        presentation_data["slides"].append(slide_data)
    output_file = output_dir / "parsed.json"
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(presentation_data, f, indent=2, ensure_ascii=False)
    return presentation_data
