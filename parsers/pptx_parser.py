"""Модуль для парсинга структуры и подробного визуального содержимого PPTX презентаций."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

from models.contracts import (
    ParsedPresentation,
    ParsedSlide,
    ParsedShape,
    ShapeStyle,
    GroupPosition
)

logger = logging.getLogger(__name__)


def _iter_all_shapes(shapes) -> list:
    """Рекурсивно собирает все элементы (включая сгруппированные) в плоский список."""
    all_shapes = []
    for shape in shapes:
        if shape.shape_type == 6:
            all_shapes.extend(_iter_all_shapes(shape.shapes))
        else:
            all_shapes.append(shape)
    return all_shapes


def _map_shape_type(shape) -> str:
    """Определяет тип шейпа в виде строки."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return "group"
        elif shape.has_text_frame:
            return "text"
        else:
            return "shape"
    except Exception:
        return "unknown"


def _extract_position(shape) -> GroupPosition:
    """Извлекает позицию и размеры шейпа."""
    try:
        return GroupPosition(
            x=shape.left / 9525 if shape.left else 0.0,
            y=shape.top / 9525 if shape.top else 0.0,
            w=shape.width / 9525 if shape.width else 0.0,
            h=shape.height / 9525 if shape.height else 0.0
        )
    except Exception:
        return GroupPosition(x=0.0, y=0.0, w=0.0, h=0.0)


def _extract_text_styles(shape) -> ShapeStyle:
    """Извлекает стили текста (шрифт, размер, цвет и начертание)."""
    style = ShapeStyle()
    try:
        if not shape.has_text_frame:
            return style
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font:
                    if run.font.name:
                        style.font_family = run.font.name
                    if run.font.size:
                        style.font_size = run.font.size.pt
                    if run.font.color and run.font.color.type == 1:
                        style.font_color = f"#{str(run.font.color.rgb)}"
                    if run.font.bold is not None:
                        style.bold = run.font.bold
                    if run.font.italic is not None:
                        style.italic = run.font.italic
                break
            break
    except Exception as e:
        logger.info(f"Ошибка при извлечении стилей текста: {e}")
    return style


def _extract_fill_lines(shape) -> ShapeStyle:
    """Извлекает цвета заливки и контура шейпа."""
    style = ShapeStyle()
    try:
        if hasattr(shape, "fill") and shape.fill.type == 1:
            if hasattr(shape.fill.fore_color, "rgb") and shape.fill.fore_color.rgb:
                style.fill_color = f"#{str(shape.fill.fore_color.rgb)}"
        if hasattr(shape, "line") and shape.line.fill.type == 1:
            if hasattr(shape.line.color, "rgb") and shape.line.color.rgb:
                style.line_color = f"#{str(shape.line.color.rgb)}"
    except Exception as e:
        logger.info(f"Ошибка при извлечении заливки/контура: {e}")
    return style


def _extract_table_data(shape) -> list[list[str]]:
    """Извлекает текстовые данные из ячеек таблицы."""
    data = []
    try:
        if shape.has_table:
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text_frame.text.strip() if cell.text_frame else "")
                data.append(row_data)
    except Exception as e:
        logger.info(f"Ошибка при извлечении данных таблицы: {e}")
    return data


def _parse_slide_shapes(slide) -> list[ParsedShape]:
    """Разбирает все шейпы на слайде в модели ParsedShape."""
    parsed_shapes = []
    for shape in _iter_all_shapes(slide.shapes):
        try:
            position = _extract_position(shape)
            shape_type = _map_shape_type(shape)
            texts = []
            if shape.has_text_frame:
                texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            
            style = _extract_text_styles(shape)
            fill_line_style = _extract_fill_lines(shape)
            style.fill_color = fill_line_style.fill_color
            style.line_color = fill_line_style.line_color
            
            table_data = _extract_table_data(shape) if shape.has_table else None
            
            parsed_shape = ParsedShape(
                shape_id=str(shape.shape_id) if hasattr(shape, "shape_id") else "",
                shape_type=shape_type,
                name=shape.name if hasattr(shape, "name") else "",
                position=position,
                rotation=shape.rotation if hasattr(shape, "rotation") and shape.rotation else 0.0,
                texts=texts,
                style=style,
                table_data=table_data,
                image_index=None
            )
            parsed_shapes.append(parsed_shape)
        except Exception as e:
            logger.info(f"Ошибка при обработке шейпа: {e}")
            continue
    return parsed_shapes


def parse_pptx_rich(pptx_path: str | Path) -> ParsedPresentation:
    """Извлекает расширенную структуру презентации со стилями и позициями."""
    path = Path(pptx_path)
    logger.info(f"Начало глубокого парсинга: {path.name}")

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.error(f"Не удалось открыть файл {pptx_path}: {e}")
        raise ValueError(f"Invalid PPTX file: {e}")

    slides_data = []
    
    for idx, slide in enumerate(prs.slides):
        bg_color = None
        try:
            if getattr(slide, "background", None) and getattr(slide.background, "fill", None) and slide.background.fill.type is not None:
                if hasattr(slide.background.fill.fore_color, "rgb") and slide.background.fill.fore_color.rgb:
                    bg_color = f"#{str(slide.background.fill.fore_color.rgb)}"
        except Exception:
            pass

        slides_data.append(ParsedSlide(
            slide_index=idx,
            width=prs.slide_width / 9525,
            height=prs.slide_height / 9525,
            background_color=bg_color,
            shapes=_parse_slide_shapes(slide)
        ))

    logger.info(f"Успешно обработано {len(slides_data)} слайдов (rich)")
    return ParsedPresentation(filename=path.name, slide_count=len(slides_data), slides=slides_data)


if __name__ == "__main__":
    import sys
    
    test_path = sys.argv[1] if len(sys.argv) > 1 else "projects/test_map/test_maps.pptx"
    
    try:
        test_result = parse_pptx_rich(test_path)
        for s in test_result.slides:
            logger.info(f"Slide {s.slide_index} | BG: {s.background_color} | Shapes: {len(s.shapes)}")
            for sh in s.shapes[:2]:
                logger.info(f"  - {sh.shape_type}: {sh.name} at {sh.position}")
    except Exception as e:
        logger.error(f"Ошибка при тестировании: {e}")