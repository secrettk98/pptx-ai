from pptx import Presentation
from lxml import etree

prs = Presentation("temp/debug.pptx")
slide = prs.slides[0]
for shape in slide.shapes:
    sp = shape._element
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    prstGeom = sp.find('.//a:prstGeom', nsmap)
    prst_val = prstGeom.get('prst') if prstGeom is not None else "NONE"
    print(f"{shape.name:20s} | prst={prst_val:15s} | type={shape.shape_type}")