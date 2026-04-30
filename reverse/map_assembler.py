"""Модуль финальной сборки карты в PPTX с заменой подложки и пересчетом позиций объектов."""

import logging
from pathlib import Path
from typing import Optional, Union, Tuple
from copy import deepcopy

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel

from reverse.map_layer_splitter import split_slide_layers, ShapeInfo
from reverse.map_background import replace_background, GeoInfo
from reverse.map_objects_redesign import redesign_objects, ObjectsRedesignResult

logger = logging.getLogger(__name__)

DEFAULT_OUTPUT_DIR = "projects/test_map/output"
MAPBOX_STYLE = "mapbox"


class AssemblyResult(BaseModel):
    """Результат сборки финального слайда."""
    output_path: Optional[str] = None
    slide_number: int
    objects_total: int
    objects_moved: int
    objects_out_of_bounds: int
    success: bool = True
    error: Optional[str] = None


def _find_background_shape(slide, bg_shape_info: ShapeInfo):
    """Ищет шейп подложки на слайде по имени или максимальной площади."""
    for shape in slide.shapes:
        if shape.name == bg_shape_info.name:
            return shape
            
    largest_pic = None
    max_area = 0
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            area = shape.width * shape.height
            if area > max_area:
                max_area = area
                largest_pic = shape
                
    return largest_pic


def _replace_bg_on_slide(slide, old_shape, new_img_path: Path, bg_info: ShapeInfo) -> Tuple[int, int, int, int]:
    """Заменяет картинку на заднем плане слайда и возвращает ее координаты."""
    bg_left = old_shape.left
    bg_top = old_shape.top
    bg_width = old_shape.width
    bg_height = old_shape.height

    sp = old_shape._element
    sp.getparent().remove(sp)

    new_bg_shape = slide.shapes.add_picture(
        str(new_img_path), bg_left, bg_top, bg_width, bg_height
    )

    sp_tree = slide.shapes._spTree
    sp_tree.remove(new_bg_shape._element)
    sp_tree.insert(2, new_bg_shape._element) 
    
    return bg_left, bg_top, bg_width, bg_height


def _relocate_objects_on_slide(slide, redesigned: list, b_left: int, b_top: int, b_width: int, b_height: int) -> int:
    """Перемещает объекты на новые позиции на слайде."""
    moved_count = 0
    for obj in redesigned:
        if obj.out_of_bounds:
            continue
            
        for shape in slide.shapes:
            if shape.name == obj.name:
                shape.left = b_left + int(obj.rel_x * b_width)
                shape.top = b_top + int(obj.rel_y * b_height)
                shape.width = int(obj.rel_width * b_width)
                shape.height = int(obj.rel_height * b_height)
                moved_count += 1
                break
                
    return moved_count


def assemble_map_slide(
    pptx_path: Union[str, Path],
    slide_number: int = 1,
    output_dir: Union[str, Path] = DEFAULT_OUTPUT_DIR,
    accent_color: str = "#0066CC",
) -> AssemblyResult:
    """Главная функция сборки слайда с новой подложкой и перемещенными объектами."""
    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        logger.info(f"Начало сборки слайда {slide_number} из {pptx_path}")
        layers = split_slide_layers(str(pptx_path), slide_number, str(layers_dir))
        
        if not layers.background or not layers.background_image_path:
            raise ValueError("Слой подложки не найден на исходном слайде")
            
        logger.info("Замена подложки")
        bg_result = replace_background(str(layers.background_image_path), MAPBOX_STYLE, str(out_dir))
        if not bg_result.success:
            raise RuntimeError("Замена подложки завершилась ошибкой")
            
        logger.info("Определение размеров новой подложки")
        with Image.open(bg_result.new_image_path) as img:
            new_w, new_h = img.size
            
        logger.info("Пересчет координат объектов")
        redesign_res = redesign_objects(
            layers.objects, layers.background[0], new_w, new_h, bg_result.geo_info
        )
        
        logger.info("Открытие презентации для модификации")
        prs = Presentation(str(pptx_path))
        slide = prs.slides[slide_number - 1]
        
        old_bg_shape = _find_background_shape(slide, layers.background[0])
        if not old_bg_shape:
            raise ValueError("Шейп старой подложки не найден в дереве элементов")
            
        logger.info("Удаление старой подложки и вставка новой")
        b_left, b_top, b_width, b_height = _replace_bg_on_slide(
            slide, old_bg_shape, Path(bg_result.new_image_path), layers.background[0]
        )
        
        logger.info("Пересадка объектов на слайде")
        moved = _relocate_objects_on_slide(
            slide, redesign_res.redesigned, b_left, b_top, b_width, b_height
        )
        
        out_path = out_dir / f"assembled_slide_{slide_number}.pptx"
        logger.info(f"Сохранение готовой презентации в {out_path}")
        prs.save(str(out_path))
        
        return AssemblyResult(
            output_path=str(out_path),
            slide_number=slide_number,
            objects_total=redesign_res.total_objects,
            objects_moved=moved,
            objects_out_of_bounds=redesign_res.out_of_bounds_count
        )

    except Exception as e:
        logger.error(f"Ошибка в процессе финальной сборки слайда: {e}")
        return AssemblyResult(
            slide_number=slide_number,
            objects_total=0,
            objects_moved=0,
            objects_out_of_bounds=0,
            success=False,
            error=str(e)
        )


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
    
    test_pptx_path = "projects/test_map/test_maps.pptx"
    result = assemble_map_slide(test_pptx_path, slide_number=1)
    
    print(f"Успех: {result.success}")
    if result.output_path:
        print(f"Сохранено: {result.output_path}")
    if result.error:
        print(f"Ошибка: {result.error}")