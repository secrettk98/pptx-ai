"""Модуль для разделения объектов слайда на фоновую подложку и плоский список объектов."""

import os
import sys
import logging
from pathlib import Path
from typing import Optional, List, Iterable, Any

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

SHAPE_TYPE_MAP = {
    MSO_SHAPE_TYPE.PICTURE: "picture",
    MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
    MSO_SHAPE_TYPE.FREEFORM: "freeform",
    MSO_SHAPE_TYPE.GROUP: "group",
    MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
    MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
}


class ShapeInfo(BaseModel):
    """Информация об одном объекте на слайде."""
    name: str
    shape_type: str
    left: float
    top: float
    width: float
    height: float
    area: float
    text: Optional[str] = None
    has_fill: bool = False
    # Crop offsets (0.0 — 1.0, доля обрезки с каждой стороны)
    crop_left: float = 0.0
    crop_right: float = 0.0
    crop_top: float = 0.0
    crop_bottom: float = 0.0


class LayerSplit(BaseModel):
    """Результат разделения слайда на слои."""
    slide_number: int
    total_shapes: int
    background: List[ShapeInfo]
    objects: List[ShapeInfo]
    background_image_path: Optional[str] = None


def _iter_all_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    """Рекурсивно обходит коллекцию шейпов, разворачивая группы."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)
        else:
            yield shape


def extract_background_image(shape: Any, output_dir: Path, slide_number: int, 
                             crop_left: float = 0.0, crop_right: float = 0.0,
                             crop_top: float = 0.0, crop_bottom: float = 0.0) -> Optional[str]:
    """Извлекает бинарные данные картинки и сохраняет в файл с учетом обрезки."""
    try:
        blob = shape.image.blob
        ext = shape.image.content_type.split("/")[-1].replace("jpeg", "jpg")
        path = output_dir / f"background_slide_{slide_number}.{ext}"
        with open(path, "wb") as f:
            f.write(blob)
            
        if crop_left > 0 or crop_right > 0 or crop_top > 0 or crop_bottom > 0:
            img = Image.open(path)
            w, h = img.size
            left_px = int(crop_left * w)
            top_px = int(crop_top * h)
            right_px = int(w - crop_right * w)
            bottom_px = int(h - crop_bottom * h)
            
            cropped = img.crop((left_px, top_px, right_px, bottom_px))
            cropped.save(path)
            logger.info(f"Подложка обрезана: {w}x{h} -> {cropped.size}")
            
        return str(path)
    except Exception as e:
        logger.error(f"Не удалось извлечь подложку со слайда {slide_number}: {e}")
        return None


def split_slide_layers(pptx_path: str | Path, slide_number: int = 1, output_dir: str | Path | None = None) -> LayerSplit:
    """Разделяет объекты слайда на фоновую картинку и остальные элементы."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_number - 1]

    slide_area = (prs.slide_width / 914400) * (prs.slide_height / 914400)
    shapes_info: List[tuple[Any, ShapeInfo]] = []

    for shape in _iter_all_shapes(slide.shapes):
        try:
            w, h = shape.width / 914400, shape.height / 914400
            has_fill = False
            try:
                if shape.fill.type is not None:
                    has_fill = True
            except Exception:
                pass

            crop_l = crop_r = crop_t = crop_b = 0.0
            try:
                blip_fill = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
                if blip_fill is not None:
                    crop_l = int(blip_fill.get('l', '0')) / 100000
                    crop_r = int(blip_fill.get('r', '0')) / 100000
                    crop_t = int(blip_fill.get('t', '0')) / 100000
                    crop_b = int(blip_fill.get('b', '0')) / 100000
            except Exception:
                pass

            info = ShapeInfo(
                name=shape.name,
                shape_type=SHAPE_TYPE_MAP.get(shape.shape_type, "other"),
                left=shape.left / 914400,
                top=shape.top / 914400,
                width=w,
                height=h,
                area=w * h,
                text=shape.text_frame.text if shape.has_text_frame else None,
                has_fill=has_fill,
                crop_left=crop_l,
                crop_right=crop_r,
                crop_top=crop_t,
                crop_bottom=crop_b,
            )

            shapes_info.append((shape, info))
        except Exception as e:
            logger.error(f"Ошибка при обработке шейпа '{getattr(shape, 'name', 'unnamed')}' на слайде {slide_number}: {e}")

    pictures = [(s, i) for s, i in shapes_info if i.shape_type == "picture"]
    bg_shape_tuple = max(pictures, key=lambda x: x[1].area, default=None)

    background = []
    objects = []
    bg_path = None

    if bg_shape_tuple and bg_shape_tuple[1].area > (0.3 * slide_area):
        background.append(bg_shape_tuple[1])
        if output_dir:
            bg_info = bg_shape_tuple[1]
            bg_path = extract_background_image(
                bg_shape_tuple[0], Path(output_dir), slide_number,
                crop_left=bg_info.crop_left,
                crop_right=bg_info.crop_right,
                crop_top=bg_info.crop_top,
                crop_bottom=bg_info.crop_bottom,
            )

    for shape, info in shapes_info:
        if bg_shape_tuple and shape == bg_shape_tuple[0]:
            continue
        objects.append(info)

    return LayerSplit(
        slide_number=slide_number,
        total_shapes=len(shapes_info),
        background=background,
        objects=objects,
        background_image_path=bg_path
    )


if __name__ == "__main__":
    pptx_f = sys.argv[1] if len(sys.argv) > 1 else "projects/test_map/test_maps.pptx"
    slide_n = int(sys.argv[2]) if len(sys.argv) > 2 else 1
    out_dir = Path("projects/test_map/layers")
    out_dir.mkdir(parents=True, exist_ok=True)

    result = split_slide_layers(pptx_f, slide_n, out_dir)
    print(f"Слайд {result.slide_number}: {result.total_shapes} объектов")
    print(f"Подложка: {len(result.background)} объектов")
    print(f"Объекты: {len(result.objects)} объектов")
    for obj in result.objects:
        print(f"  - {obj.name} ({obj.shape_type}) {obj.width:.1f}x{obj.height:.1f} inches")
    if result.background_image_path:
        print(f"Подложка сохранена: {result.background_image_path}")