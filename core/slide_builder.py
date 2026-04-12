"""
Module for building PowerPoint presentations using python-pptx.
"""

from pptx.util import Inches, Pt
import logging

logger = logging.getLogger(__name__)


def add_text_box(slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 18) -> None:
    """
    Add a text box to a slide with specified dimensions in inches.
    
    Args:
        slide: The slide object to add the text box to
        text (str): The text content for the text box
        left (float): The left position of the text box in inches
        top (float): The top position of the text box in inches
        width (float): The width of the text box in inches
        height (float): The height of the text box in inches
        font_size (int): The font size for the text (default: 18)
    """
    try:
        left_emu = Inches(left)
        top_emu = Inches(top)
        width_emu = Inches(width)
        height_emu = Inches(height)
        
        text_box = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
        text_frame = text_box.text_frame
        text_frame.text = text
        
        if text_frame.paragraphs:
            text_frame.paragraphs[0].font.size = Pt(font_size)
   	
        logger.info(f"Text box added successfully with dimensions {width}x{height} inches")
        
    except Exception as e:
        logger.error(f"Failed to add text box to slide: {str(e)}")
        raise