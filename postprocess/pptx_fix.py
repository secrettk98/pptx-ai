"""Модуль постобработки и форматирования сгенерированных PPTX-файлов."""

import sys
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

from core.utils.logger import get_logger

log = get_logger(__name__)

INNER_MARGIN = Pt(0)
FONT_FAMILY = "Google Sans"


def fix_text_margins(shape: Any) -> bool:
    """Устанавливает нулевые внутренние отступы текстового фрейма."""
    try:
        if not getattr(shape, "has_text_frame", False):
            return False
            
        tf = shape.text_frame
        tf.margin_left = INNER_MARGIN
        tf.margin_right = INNER_MARGIN
        tf.margin_top = INNER_MARGIN
        tf.margin_bottom = INNER_MARGIN
        tf.word_wrap = True
        
        return True
    except Exception as e:
        log.error(f"Ошибка при настройке отступов текста: {e}")
        raise


def fix_fonts(shape: Any) -> bool:
    """Устанавливает целевой шрифт для всех текстовых элементов без явно заданного шрифта."""
    try:
        if not getattr(shape, "has_text_frame", False):
            return False
            
        changed = False
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.name is None or run.font.name == "":
                    run.font.name = FONT_FAMILY
                    changed = True
                    
        return changed
    except Exception as e:
        log.error(f"Ошибка при замене шрифтов: {e}")
        raise


def fix_rounded_corners(shape):
    """Rectangle -> roundRect. TextBox -> обычный rect (без скруглений)."""
    name = shape.name or ""
    sp = shape._element
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    prstGeom = sp.find('.//a:prstGeom', nsmap)
    
    if prstGeom is None:
        return False
    
    prst = prstGeom.get('prst', '')
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    
    # TextBox — убрать скругление, сделать обычный rect
    if name.startswith("TextBox"):
        if prst == 'roundRect':
            prstGeom.set('prst', 'rect')
            old_avLst = prstGeom.find(f'{{{ns}}}avLst')
            if old_avLst is not None:
                prstGeom.remove(old_avLst)
            etree.SubElement(prstGeom, f'{{{ns}}}avLst')
            return True
        return False
    
    # Rectangle — добавить мягкое скругление
    if name.startswith("Rectangle"):
        if prst == 'rect' or prst == 'roundRect':
            prstGeom.set('prst', 'roundRect')
            old_avLst = prstGeom.find(f'{{{ns}}}avLst')
            if old_avLst is not None:
                prstGeom.remove(old_avLst)
            avLst = etree.SubElement(prstGeom, f'{{{ns}}}avLst')
            gd = etree.SubElement(avLst, f'{{{ns}}}gd')
            gd.set('name', 'adj')
            gd.set('fmla', 'val 5000')
            return True
    
    return False


def _iter_all_shapes(shapes: Any) -> list[Any]:
    """Рекурсивно извлекает все фигуры на слайде, включая вложенные в группы."""
    try:
        result = []
        for shape in shapes:
            if getattr(shape, "shape_type", None) == 6:
                result.extend(_iter_all_shapes(shape.shapes))
            else:
                result.append(shape)
        return result
    except Exception as e:
        log.error(f"Ошибка при рекурсивном обходе фигур: {e}")
        raise


def fix_pptx(pptx_path: str | Path) -> Path:
    """Открывает PPTX-файл, применяет исправления и сохраняет изменения."""
    try:
        pptx_path = Path(pptx_path)
        log.info(f"PPTX post-processing: {pptx_path.name}")
        
        prs = Presentation(str(pptx_path))
        margins_fixed = 0
        fonts_fixed = 0
        corners_fixed = 0
        
        for slide in prs.slides:
            for shape in _iter_all_shapes(slide.shapes):
                if fix_text_margins(shape):
                    margins_fixed += 1
                if fix_fonts(shape):
                    fonts_fixed += 1
                if fix_rounded_corners(shape):
                    corners_fixed += 1
                    
        prs.save(str(pptx_path))
        log.info(f"  Fixed: {margins_fixed} margins, {fonts_fixed} fonts, {corners_fixed} corners")
        
        return pptx_path
    except Exception as e:
        log.error(f"Ошибка при постобработке файла {pptx_path}: {e}")
        raise


if __name__ == "__main__":
    try:
        path = sys.argv[1] if len(sys.argv) > 1 else "output/redesigned_test_maps.pptx"
        fix_pptx(path)
        print("Done!")
    except Exception as e:
        log.error(f"Критическая ошибка выполнения скрипта: {e}")
        sys.exit(1)